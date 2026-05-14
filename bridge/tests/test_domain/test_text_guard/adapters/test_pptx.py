"""Tests for the PPTX adapter (via TextGuardService).

These tests are skipped if python-pptx is not installed.
"""

from __future__ import annotations

from pathlib import Path

import pytest

try:
    from pptx import Presentation

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from application.text_guard_service import TextGuardService

pytestmark = pytest.mark.skipif(not HAS_PPTX, reason="python-pptx not installed")


@pytest.fixture
def tg_service() -> TextGuardService:
    """Text Guard service instance."""
    return TextGuardService()


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    """Create a sample PPTX with ASCII umlauts."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Erklaerung fuer das Projekt"

    subtitle = slide.placeholders[1]
    subtitle.text = "Natuerlich ist das moeglich"

    filepath = tmp_path / "test.pptx"
    prs.save(str(filepath))
    return filepath


@pytest.fixture
def clean_pptx(tmp_path: Path) -> Path:
    """Create a PPTX with correct umlauts."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Erklärung für das Projekt"
    filepath = tmp_path / "clean.pptx"
    prs.save(str(filepath))
    return filepath


class TestPptxCheck:
    """Tests for checking PPTX files."""

    def test_finds_issues(
        self, tg_service: TextGuardService, sample_pptx: Path
    ) -> None:
        """Detects ASCII umlauts in PPTX text frames."""
        issues = tg_service.check_file(sample_pptx, "de")
        assert len(issues) >= 2

    def test_clean_file(self, tg_service: TextGuardService, clean_pptx: Path) -> None:
        """Clean PPTX returns no issues."""
        issues = tg_service.check_file(clean_pptx, "de")
        assert issues == []


class TestPptxFix:
    """Tests for fixing PPTX files."""

    def test_fixes_in_place(
        self, tg_service: TextGuardService, sample_pptx: Path
    ) -> None:
        """Fixes ASCII umlauts in PPTX in-place."""
        changed = tg_service.fix_file(sample_pptx, "de")
        assert changed is True

        # Verify the fix
        prs = Presentation(str(sample_pptx))
        slide = prs.slides[0]
        title_text = slide.shapes.title.text
        assert "für" in title_text or "Erklärung" in title_text

    def test_no_change_on_clean(
        self, tg_service: TextGuardService, clean_pptx: Path
    ) -> None:
        """Clean PPTX is not modified."""
        changed = tg_service.fix_file(clean_pptx, "de")
        assert changed is False
