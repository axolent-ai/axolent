"""Text Guard application service.

Coordinates text guard operations across different file formats.
Loads YAML rule files and creates domain guard instances.
Provides the integration point for the bot pipeline and CLI.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Literal

from domain.text_guard import TextGuard, get_builtin_rules
from domain.text_guard.adapters.streaming import StreamingTextGuard
from domain.text_guard.adapters.text import check_text, fix_text
from domain.text_guard.models import Issue, RuleSet, WordPair

log = logging.getLogger(__name__)

# Supported text file extensions for the text adapter
TEXT_EXTENSIONS: frozenset[str] = frozenset(
    {".md", ".txt", ".json", ".yaml", ".yml", ".html", ".htm", ".xml", ".csv"}
)

# Office file extensions (require optional dependencies)
OFFICE_EXTENSIONS: frozenset[str] = frozenset({".docx", ".xlsx", ".pptx"})


class TextGuardService:
    """Application service for text guard operations.

    Creates and caches TextGuard instances per language.
    Provides file-level operations (check/fix files).
    """

    def __init__(self) -> None:
        self._guards: dict[str, TextGuard] = {}
        self._streaming_guards: dict[str, StreamingTextGuard] = {}

    def get_guard(
        self, language: str, *, mode: Literal["check", "fix"] = "fix"
    ) -> TextGuard | None:
        """Get or create a TextGuard for a language.

        Args:
            language: ISO 639-1 code.
            mode: "check" or "fix".

        Returns:
            TextGuard instance, or None if no rules exist.
        """
        cache_key = f"{language}:{mode}"
        if cache_key not in self._guards:
            rules = get_builtin_rules(language)
            if rules is None:
                return None
            self._guards[cache_key] = TextGuard(rules, mode=mode)
        return self._guards[cache_key]

    def get_streaming_guard(self, language: str) -> StreamingTextGuard | None:
        """Get or create a StreamingTextGuard for a language.

        Args:
            language: ISO 639-1 code.

        Returns:
            StreamingTextGuard instance, or None if no rules exist.
        """
        if language not in self._streaming_guards:
            guard = self.get_guard(language, mode="fix")
            if guard is None:
                return None
            self._streaming_guards[language] = StreamingTextGuard(guard)
        return self._streaming_guards[language]

    def check_string(self, text: str, language: str) -> list[Issue]:
        """Check a string for diacritic issues.

        Args:
            text: Text to check.
            language: ISO 639-1 code.

        Returns:
            List of issues, empty if language not supported.
        """
        guard = self.get_guard(language, mode="check")
        if guard is None:
            return []
        return check_text(text, guard)

    def fix_string(self, text: str, language: str) -> str:
        """Fix diacritic issues in a string.

        Args:
            text: Text to fix.
            language: ISO 639-1 code.

        Returns:
            Corrected text (unchanged if language not supported).
        """
        guard = self.get_guard(language, mode="fix")
        if guard is None:
            return text
        return fix_text(text, guard)

    def check_file(self, filepath: Path, language: str) -> list[Issue]:
        """Check a file for diacritic issues.

        Dispatches to the appropriate adapter based on file extension.

        Args:
            filepath: Path to the file.
            language: ISO 639-1 code.

        Returns:
            List of issues found.
        """
        guard = self.get_guard(language, mode="check")
        if guard is None:
            return []

        ext = filepath.suffix.lower()

        if ext in TEXT_EXTENSIONS:
            content = filepath.read_text(encoding="utf-8", errors="replace")
            return check_text(content, guard)

        if ext == ".docx":
            return self._check_docx(filepath, guard)
        if ext == ".xlsx":
            return self._check_xlsx(filepath, guard)
        if ext == ".pptx":
            return self._check_pptx(filepath, guard)

        log.warning("Unsupported file type: %s", ext)
        return []

    def fix_file(self, filepath: Path, language: str) -> bool:
        """Fix diacritic issues in a file in-place.

        Args:
            filepath: Path to the file.
            language: ISO 639-1 code.

        Returns:
            True if changes were made, False otherwise.
        """
        guard = self.get_guard(language, mode="fix")
        if guard is None:
            return False

        ext = filepath.suffix.lower()

        if ext in TEXT_EXTENSIONS:
            return self._fix_text_file(filepath, guard)

        if ext == ".docx":
            return self._fix_docx(filepath, guard)
        if ext == ".xlsx":
            return self._fix_xlsx(filepath, guard)
        if ext == ".pptx":
            return self._fix_pptx(filepath, guard)

        log.warning("Unsupported file type: %s", ext)
        return False

    @staticmethod
    def _fix_text_file(filepath: Path, guard: TextGuard) -> bool:
        """Fix a plain text file in-place."""
        content = filepath.read_text(encoding="utf-8", errors="replace")
        fixed = fix_text(content, guard)
        if fixed != content:
            filepath.write_text(fixed, encoding="utf-8")
            return True
        return False

    @staticmethod
    def _check_docx(filepath: Path, guard: TextGuard) -> list[Issue]:
        """Check a DOCX file for issues."""
        try:
            from docx import Document
        except ImportError:
            log.warning(
                "python-docx not installed. Install with: pip install python-docx"
            )
            return []

        doc = Document(str(filepath))
        all_issues: list[Issue] = []
        for para in doc.paragraphs:
            issues = guard.check(para.text)
            all_issues.extend(issues)
        return all_issues

    @staticmethod
    def _fix_docx(filepath: Path, guard: TextGuard) -> bool:
        """Fix a DOCX file in-place."""
        try:
            from docx import Document
        except ImportError:
            log.warning(
                "python-docx not installed. Install with: pip install python-docx"
            )
            return False

        doc = Document(str(filepath))
        changed = False
        for para in doc.paragraphs:
            for run in para.runs:
                fixed = guard.fix(run.text)
                if fixed != run.text:
                    run.text = fixed
                    changed = True
        if changed:
            doc.save(str(filepath))
        return changed

    @staticmethod
    def _check_xlsx(filepath: Path, guard: TextGuard) -> list[Issue]:
        """Check an XLSX file for issues (skips formula cells)."""
        try:
            from openpyxl import load_workbook
        except ImportError:
            log.warning("openpyxl not installed. Install with: pip install openpyxl")
            return []

        wb = load_workbook(str(filepath), data_only=False)
        all_issues: list[Issue] = []
        for ws in wb.worksheets:
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value is None:
                        continue
                    if isinstance(cell.value, str):
                        # Skip formula cells
                        if cell.value.startswith("="):
                            continue
                        issues = guard.check(cell.value)
                        all_issues.extend(issues)
        wb.close()
        return all_issues

    @staticmethod
    def _fix_xlsx(filepath: Path, guard: TextGuard) -> bool:
        """Fix an XLSX file in-place (skips formula cells)."""
        try:
            from openpyxl import load_workbook
        except ImportError:
            log.warning("openpyxl not installed. Install with: pip install openpyxl")
            return False

        wb = load_workbook(str(filepath), data_only=False)
        changed = False
        for ws in wb.worksheets:
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value is None:
                        continue
                    if isinstance(cell.value, str):
                        # Skip formula cells
                        if cell.value.startswith("="):
                            continue
                        fixed = guard.fix(cell.value)
                        if fixed != cell.value:
                            cell.value = fixed
                            changed = True
        if changed:
            wb.save(str(filepath))
        wb.close()
        return changed

    @staticmethod
    def _check_pptx(filepath: Path, guard: TextGuard) -> list[Issue]:
        """Check a PPTX file for issues."""
        try:
            from pptx import Presentation
        except ImportError:
            log.warning(
                "python-pptx not installed. Install with: pip install python-pptx"
            )
            return []

        prs = Presentation(str(filepath))
        all_issues: list[Issue] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        issues = guard.check(run.text)
                        all_issues.extend(issues)
        return all_issues

    @staticmethod
    def _fix_pptx(filepath: Path, guard: TextGuard) -> bool:
        """Fix a PPTX file in-place."""
        try:
            from pptx import Presentation
        except ImportError:
            log.warning(
                "python-pptx not installed. Install with: pip install python-pptx"
            )
            return False

        prs = Presentation(str(filepath))
        changed = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        fixed = guard.fix(run.text)
                        if fixed != run.text:
                            run.text = fixed
                            changed = True
        if changed:
            prs.save(str(filepath))
        return changed


def load_rules_from_yaml(yaml_path: Path) -> RuleSet | None:
    """Load a rule set from a YAML file.

    Allows extending built-in rules with custom YAML files.

    Args:
        yaml_path: Path to the YAML rule file.

    Returns:
        RuleSet if successfully loaded, None on error.
    """
    try:
        import yaml
    except ImportError:
        log.warning("PyYAML not installed. Cannot load YAML rules.")
        return None

    try:
        content = yaml_path.read_text(encoding="utf-8")
        data = yaml.safe_load(content)
    except Exception as exc:
        log.error("Failed to load YAML rules from %s: %s", yaml_path, exc)
        return None

    if not data or not isinstance(data, dict):
        return None

    language = data.get("language", "")
    pairs_data = data.get("word_pairs", [])
    whitelist_data = data.get("loan_word_whitelist", [])

    word_pairs = tuple(
        WordPair(
            ascii_form=p["ascii"],
            correct_form=p["correct"],
            case_sensitive=p.get("case_sensitive", False),
        )
        for p in pairs_data
        if isinstance(p, dict) and "ascii" in p and "correct" in p
    )

    whitelist = frozenset(str(w) for w in whitelist_data)

    return RuleSet(
        language=language,
        word_pairs=word_pairs,
        loan_word_whitelist=whitelist,
    )
